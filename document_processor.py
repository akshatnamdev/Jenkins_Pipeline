"""Document processing module for DRAVIS"""
import json
import hashlib
from datetime import datetime
from pathlib import Path
from typing import List, Dict

import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import pytesseract

from config import settings


class DocumentProcessor:
    """Handles document parsing and processing"""

    SUPPORTED_FORMATS = {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "txt": "text/plain",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "png": "image/png",
        "bmp": "image/bmp",
        "py": "text/plain",
        "java": "text/plain",
        "cpp": "text/plain",
        "js": "text/plain",
        "json": "application/json",
    }

    def __init__(self):
        self.documents_dir: Path = settings.DOCUMENTS_DIR
        self.documents_dir.mkdir(parents=True, exist_ok=True)
        self.metadata_file = self.documents_dir / "metadata.json"
        self.metadata: Dict[str, Dict] = self._load_metadata()

    # ------------------------------------------------------------------ #
    # METADATA
    # ------------------------------------------------------------------ #
    def _load_metadata(self) -> Dict:
        if self.metadata_file.exists():
            with open(self.metadata_file, "r", encoding="utf-8") as f:
                return json.load(f)
        return {}

    def _save_metadata(self):
        with open(self.metadata_file, "w", encoding="utf-8") as f:
            json.dump(self.metadata, f, indent=2, default=str, ensure_ascii=False)

    # ------------------------------------------------------------------ #
    # RAW TEXT EXTRACTORS
    # ------------------------------------------------------------------ #
    def process_pdf(self, file_path: Path) -> str:
        text_parts = []
        try:
            with open(file_path, "rb") as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text() or ""
                    text_parts.append(f"--- Page {page_num + 1} ---\n{page_text}")
        except Exception as e:
            print(f"Error reading PDF {file_path}: {e}")
            return ""
        return "\n\n".join(text_parts)

    def process_docx(self, file_path: Path) -> str:
        try:
            doc = DocxDocument(file_path)
            paragraphs = [para.text for para in doc.paragraphs]
            return "\n".join(paragraphs)
        except Exception as e:
            print(f"Error reading DOCX {file_path}: {e}")
            return ""

    def process_pptx(self, file_path: Path) -> str:
        try:
            prs = Presentation(file_path)
            text = []
            for slide_num, slide in enumerate(prs.slides):
                text.append(f"--- Slide {slide_num + 1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            print(f"Error reading PPTX {file_path}: {e}")
            return ""

    def process_image(self, file_path: Path) -> str:
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text if text.strip() else f"[Image: {file_path.name}]"
        except Exception as e:
            print(f"Error reading image {file_path}: {e}")
            return f"[Image: {file_path.name}]"

    def process_text_file(self, file_path: Path) -> str:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
            with open(file_path, "r", encoding="latin-1") as f:
                return f.read()
        except Exception as e:
            print(f"Error reading text file {file_path}: {e}")
            return ""

    def extract_text(self, file_path: Path) -> str:
        ext = file_path.suffix.lower().strip(".")
        if ext == "pdf":
            return self.process_pdf(file_path)
        if ext == "docx":
            return self.process_docx(file_path)
        if ext == "pptx":
            return self.process_pptx(file_path)
        if ext in ["jpg", "jpeg", "png", "bmp"]:
            return self.process_image(file_path)
        if ext in ["txt", "py", "java", "cpp", "js", "json"]:
            return self.process_text_file(file_path)
        return ""

    # ------------------------------------------------------------------ #
    # CHUNKING + DOCUMENT REGISTRY
    # ------------------------------------------------------------------ #
    def chunk_text(self, text: str, chunk_size: int = 512, overlap: int = 50) -> List[str]:
        words = text.split()
        chunks: List[str] = []
        for i in range(0, len(words), max(chunk_size - overlap, 1)):
            chunk = " ".join(words[i : i + chunk_size])
            if chunk.strip():
                chunks.append(chunk)
        return chunks

    def _create_doc_id(self, original_filename: str) -> str:
        """Generate a stable but unique doc id (filename + timestamp)."""
        base = f"{original_filename}-{datetime.now().isoformat()}"
        return hashlib.md5(base.encode()).hexdigest()[:10]

    def add_document(self, file_path: Path, original_filename: str) -> Dict:
        """Add a document, create chunks, and store metadata."""
        doc_id = self._create_doc_id(original_filename)

        # Ensure file is stored under DOCUMENTS_DIR
        # (file_path might already be there; if not, copy/move)
        if file_path.parent != self.documents_dir:
            # copy into documents_dir with safe name
            target_path = self.documents_dir / original_filename
            if file_path != target_path:
                target_path.write_bytes(file_path.read_bytes())
            file_path = target_path

        text = self.extract_text(file_path)
        chunks = self.chunk_text(
            text,
            getattr(settings, "CHUNK_SIZE", 512),
            getattr(settings, "CHUNK_OVERLAP", 50),
        )

        size_mb = file_path.stat().st_size / (1024 * 1024)

        self.metadata[doc_id] = {
            "filename": original_filename,
            "file_path": str(file_path),
            "file_size_mb": round(size_mb, 2),
            "upload_time": datetime.now().isoformat(),
            "chunk_count": len(chunks),
            "text_length": len(text),
            "file_type": file_path.suffix.lower(),
        }
        self._save_metadata()

        return {
            "doc_id": doc_id,
            "filename": original_filename,
            "chunks": chunks,
            "metadata": self.metadata[doc_id],
        }

    def delete_document(self, doc_id: str) -> bool:
        if doc_id in self.metadata:
            file_path = Path(self.metadata[doc_id]["file_path"])
            if file_path.exists():
                try:
                    file_path.unlink()
                except Exception:
                    pass
            del self.metadata[doc_id]
            self._save_metadata()
            return True
        return False

    def get_all_documents(self) -> Dict:
        return self.metadata

    def get_document_metadata(self, doc_id: str) -> Dict:
        return self.metadata.get(doc_id) or {}

    def get_document_text(self, doc_id: str) -> str:
        meta = self.metadata.get(doc_id)
        if not meta:
            return ""
        path = Path(meta["file_path"])
        if not path.exists():
            return ""
        return self.extract_text(path)


document_processor = DocumentProcessor()
